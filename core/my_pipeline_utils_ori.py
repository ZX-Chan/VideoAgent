import os
import json
import yaml
from jinja2 import Environment, StrictUndefined
from camel.models import ModelFactory
from camel.agents import ChatAgent
from utils.wei_utils import account_token
from utils.src.utils import get_json_from_response
from PosterAgent.gen_outline_layout import filter_image_table
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# 复用图片/表格过滤

def run_filter_image_table(args, agent_config):
    """
    复用原有filter_image_table，返回过滤后的token消耗。
    """
    return filter_image_table(args, agent_config)

# 旁白生成agent

def generate_narration_with_agent(raw_content, agent_config=None, prompt_path=None):
    """
    调用旁白agent，输入raw_content，返回每个section一段连贯学术旁白。
    prompt_path: prompt yaml路径
    agent_config: 可选，默认4o
    """
    if agent_config is None:
        from utils.wei_utils import get_agent_config
        agent_config = get_agent_config('4o')
    if prompt_path is None:
        prompt_path = 'utils/prompt_templates/my_narration_agent.yaml'
    with open(prompt_path, 'r') as f:
        prompt_config = yaml.safe_load(f)
    jinja_env = Environment(undefined=StrictUndefined)
    template = jinja_env.from_string(prompt_config["template"])
    jinja_args = {"raw_content": raw_content}
    prompt = template.render(**jinja_args)
    model = ModelFactory.create(
        model_platform=agent_config['model_platform'],
        model_type=agent_config['model_type'],
        model_config_dict=agent_config['model_config'],
        url=agent_config.get('url', None)
    )
    agent = ChatAgent(
        system_message=prompt_config['system_prompt'],
        model=model,
        message_window_size=10,
    )
    agent.reset()
    response = agent.step(prompt)
    input_token, output_token = account_token(response)
    narration_json = get_json_from_response(response.msgs[0].content)
    
    # 添加封面页的旁白
    if raw_content and "meta" in raw_content:
        meta = raw_content.get("meta", {})
        paper_title = meta.get("poster_title", meta.get("title", "本次学术报告"))
        authors = meta.get("authors", "")
        affiliations = meta.get("affiliations", "")
        
        # 生成封面页的欢迎文本（英文）
        cover_narration = f"Welcome to today's academic presentation. We will be discussing the paper titled '{paper_title}'."
        if authors:
            cover_narration += f" The authors include: {authors}."
        if affiliations:
            cover_narration += f" From: {affiliations}."
        cover_narration += " Let's begin today's academic discussion."
        
        # 将封面页旁白添加到narration_json的开头
        # 使用固定的键名，与TTS文件中的键匹配
        narration_json = {"Poster Title & Author": cover_narration, **narration_json}
    
    return narration_json, input_token, output_token

# PPT结构生成agent

def generate_ppt_from_agent(content, images, tables, agent_config=None, prompt_path=None, template_path=None, reference_path=None):
    """
    调用PPT结构agent，输入内容和图片，返回结构化PPT描述。
    prompt_path: prompt yaml路径
    template_path: 未来支持模板渲染
    reference_path: 参考PPTX路径，用于agent学习内容排布
    agent_config: 可选，默认4o
    """
    # TODO: 后续可在此处分析reference_path并注入prompt
    if agent_config is None:
        from utils.wei_utils import get_agent_config
        agent_config = get_agent_config('4o')
    if prompt_path is None:
        prompt_path = 'utils/prompt_templates/my_ppt_layout_agent.yaml'
    with open(prompt_path, 'r') as f:
        prompt_config = yaml.safe_load(f)
    jinja_env = Environment(undefined=StrictUndefined)
    template = jinja_env.from_string(prompt_config["template"])
    jinja_args = {
        "content": content,
        "images": json.dumps(images, indent=2),
        "tables": json.dumps(tables, indent=2),
        # 未来可加reference分析结果
    }
    prompt = template.render(**jinja_args)
    model = ModelFactory.create(
        model_platform=agent_config['model_platform'],
        model_type=agent_config['model_type'],
        model_config_dict=agent_config['model_config'],
        url=agent_config.get('url', None)
    )
    agent = ChatAgent(
        system_message=prompt_config['system_prompt'],
        model=model,
        message_window_size=10,
    )
    agent.reset()
    response = agent.step(prompt)
    input_token, output_token = account_token(response)
    ppt_json = get_json_from_response(response.msgs[0].content)
    return ppt_json, input_token, output_token

# TTS: 用OpenAI API生成中文女声mp3

def synthesize_tts_with_openai(narration_json, output_dir, openai_api_key, voice="alloy", language="zh"):
    """
    narration_json: {section: text, ...}
    output_dir: mp3输出目录
    openai_api_key: 你的OpenAI API Key
    voice: OpenAI支持的voice（如"alloy"，"echo"等，中文女声推荐alloy或官方文档最新）
    language: "zh"（中文）
    """
    import openai
    os.makedirs(output_dir, exist_ok=True)
    openai_api_base_url = os.environ.get('OPENAI_API_BASE_URL', 'https://api.openai.com/v1')
    client = openai.OpenAI(api_key=openai_api_key, base_url=openai_api_base_url)
    audio_files = {}
    for idx, (section, text) in enumerate(narration_json.items(), 1):
        audio_path = os.path.join(output_dir, f"section_{idx:02d}_{section}.mp3")
        try:
            response = client.audio.speech.create(
                model="tts-1",
                voice=voice,
                input=text,
                response_format="mp3",
                speed=1.0
            )
            with open(audio_path, "wb") as f:
                f.write(response.content)
            print(f"[TTS] {section} -> {audio_path}")
            audio_files[section] = audio_path
        except Exception as e:
            print(f"[TTS ERROR] {section}: {e}")
            audio_files[section] = None
    return audio_files

# 简单PPT渲染适配：每页正文+图片

def render_pptx_from_json(ppt_json, output_path, images, tables, width=13.33, height=7.5, template_path=None):
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    for slide_info in ppt_json:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白页
        # 完全复刻原pipeline：panel->text->figure
        panel_boxes = slide_info.get('panel_boxes', [])
        text_boxes = slide_info.get('text_boxes', [])
        figure_boxes = slide_info.get('figure_boxes', [])
        # 1. 插入panel（大区块背景/分区，可选）
        for pbox in panel_boxes:
            shape = slide.shapes.add_textbox(
                Inches(pbox['x']), Inches(pbox['y']), Inches(pbox['width']), Inches(pbox['height'])
            )
            frame = shape.text_frame
            frame.text = ''  # panel本身不放内容
            for para in frame.paragraphs:
                para.font.size = Pt(40)
        # 2. 插入text（正文/标题）
        text_content = slide_info.get('text', '')
        # 区分title和正文
        for i, tbox in enumerate(text_boxes):
            shape = slide.shapes.add_textbox(
                Inches(tbox['x']), Inches(tbox['y']), Inches(tbox['width']), Inches(tbox['height'])
            )
            frame = shape.text_frame
            frame.word_wrap = True  # 关键：自动换行
            frame.clear()
            if i == 0 and 'title' in slide_info:
                p = frame.add_paragraph()
                p.text = slide_info['title']
                p.font.size = Pt(46)
                p.font.bold = True
            else:
                p = frame.add_paragraph()
                p.text = text_content
                p.font.size = Pt(28)
                p.line_spacing = 1.3
        # 3. 插入figure（图片）
        img_list = slide_info.get('images', [])
        for i, fbox in enumerate(figure_boxes):
            if i >= len(img_list):
                break
            img_key = img_list[i]
            img_info = images.get(str(img_key), None)
            if img_info and 'image_path' in img_info:
                img_path = img_info['image_path']
            else:
                print(f"[PPTX] Warning: image key {img_key} not found in images dict or missing image_path, skip.")
                continue
            if os.path.exists(img_path):
                try:
                    slide.shapes.add_picture(
                        img_path,
                        Inches(fbox['x']), Inches(fbox['y']),
                        Inches(fbox['width']), Inches(fbox['height'])
                    )
                except Exception as e:
                    print(f"[PPTX] Failed to add image {img_path}: {e}")
            else:
                print(f"[PPTX] Image file does not exist: {img_path}")
        # 兼容无布局时的老逻辑
        if not (panel_boxes or text_boxes or figure_boxes):
            # 标题
            if 'title' in slide_info:
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(width-1), Inches(1.2))
                title_frame = title_shape.text_frame
                title_frame.text = slide_info['title']
                title_frame.paragraphs[0].font.size = Pt(54)
                title_frame.paragraphs[0].font.bold = True
            # 正文
            if 'text' in slide_info:
                content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(width-1), Inches(3.5))
                content_frame = content_shape.text_frame
                content_frame.text = slide_info['text']
                content_frame.paragraphs[0].font.size = Pt(40)
                for para in content_frame.paragraphs:
                    para.line_spacing = 1.3
            img_list = slide_info.get('images', [])
            if img_list:
                for i, img_key in enumerate(img_list):
                    img_info = images.get(str(img_key), None)
                    if img_info and 'image_path' in img_info:
                        img_path = img_info['image_path']
                    else:
                        print(f"[PPTX] Warning: image key {img_key} not found in images dict or missing image_path, skip.")
                        continue
                    if os.path.exists(img_path):
                        try:
                            with Image.open(img_path) as im:
                                img_w, img_h = im.size
                            max_w = width * 0.8
                            max_h = height * 0.4
                            img_ratio = img_w / img_h
                            box_ratio = max_w / max_h
                            if img_ratio > box_ratio:
                                disp_w = max_w
                                disp_h = max_w / img_ratio
                            else:
                                disp_h = max_h
                                disp_w = max_h * img_ratio
                            left = (width - disp_w) / 2
                            top = height - disp_h - 0.5
                            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))
                        except Exception as e:
                            print(f"[PPTX] Failed to add image {img_path}: {e}")
                    else:
                        print(f"[PPTX] Image file does not exist: {img_path}")
    prs.save(output_path)
    print(f"[PPTX] Saved to {output_path}")

def render_pptx_panel_per_slide(
    panels,
    text_arrangement,
    figure_arrangement,
    bullet_content,
    images,
    tables,
    output_path,
    width=13.33,
    height=7.5,
    template_path=None
):
    """
    复用new_pipeline的布局和内容生成方式，但渲染时每个panel/section单独生成一页slide。
    panels: section/panel结构（顺序与内容一致）
    text_arrangement: 所有文本框布局
    figure_arrangement: 所有图片框布局
    bullet_content: 每个panel的内容（与panels顺序一致）
    images, tables: 过滤后的图片/表格dict
    output_path: 输出pptx路径
    width, height: ppt页面尺寸（英寸）
    template_path: 可选模板
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    # 假设text_arrangement和figure_arrangement已按panel分组
    for idx, panel in enumerate(panels):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 标题
        title = bullet_content[idx].get('title', panel.get('section_name', f'Section {idx+1}'))
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(width-1), Inches(1.0))
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        # 正文
        text = bullet_content[idx].get('text', '')
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(width-1), Inches(3.5))
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        content_frame.clear()
        p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(28)
        p.line_spacing = 1.3
        # 图片（如有）
        img_list = bullet_content[idx].get('images', [])
        for i, img_key in enumerate(img_list):
            img_info = images.get(str(img_key), None)
            if img_info and 'image_path' in img_info:
                img_path = img_info['image_path']
            else:
                continue
            if os.path.exists(img_path):
                try:
                    # 简单居中下方插入
                    disp_w = width * 0.6
                    disp_h = height * 0.3
                    left = (width - disp_w) / 2
                    top = height - disp_h - 0.5
                    slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))
                except Exception as e:
                    print(f"[PPTX] Failed to add image {img_path}: {e}")
    prs.save(output_path)
    print(f"[PPTX] Saved to {output_path}") 