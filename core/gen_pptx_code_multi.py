import re
import json
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import math
import os


def generate_equation_grid_code(slide_var, equations_list, slide_idx_suffix, start_top_inch, available_height_inch,
                                slide_width_inch):
    """
    生成用于创建公式网格布局的Python代码。
    这是一个辅助函数，旨在保持主函数逻辑的清晰。
    """

    code = []

    # 定义布局参数
    left_margin = 1.0
    element_height = 0.5  # 文本框和图片的高度 (英寸)
    element_width = 5.0
    vertical_gap = 0.2  # 垂直间距
    horizontal_gap = 0.5  # 水平列间距
    label_textbox_width = 5.0  # 公式标签文本框的宽度

    # 根据可用高度计算一列最多能放多少行
    total_element_height_with_gap = element_height + vertical_gap
    if total_element_height_with_gap > 0:
        max_rows = max(1, math.floor(available_height_inch / total_element_height_with_gap))
    else:
        max_rows = 1

    # 为确保变量名唯一，使用 slide_idx_suffix
    sfx = slide_idx_suffix

    # 生成用于在PPT脚本中进行布局计算的代码
    code.append(f"\n# --- 开始为幻灯片(后缀: {sfx})生成公式网格布局 ---")
    code.append(f"top_margin_inch_{sfx} = {start_top_inch}")
    code.append(f"element_height_inch_{sfx} = {element_height}")
    code.append(f"element_width_inch_{sfx} = {element_width}")
    code.append(f"vertical_gap_inch_{sfx} = {vertical_gap}")
    code.append(f"horizontal_gap_inch_{sfx} = {horizontal_gap}")
    code.append(f"label_textbox_width_inch_{sfx} = {label_textbox_width}")
    code.append(f"max_rows_{sfx} = {max_rows}")
    code.append(f"current_col_x_{sfx} = Inches({left_margin})")
    code.append(f"max_width_in_col_{sfx} = 0")
    code.append(f"row_index_{sfx} = 0")

    # 遍历公式，生成放置代码
    for eq_idx, equation in enumerate(equations_list):
        equ_img_path = equation.get("equ_path", "").replace('\\', '/')
        equ_text_label = equation.get("context", "")

        code.append(f"\n# 放置公式 {eq_idx + 1} (后缀: {sfx})")

        # 计算Y坐标
        code.append(
            f"y_pos = Inches(top_margin_inch_{sfx} + row_index_{sfx} * (element_height_inch_{sfx} + vertical_gap_inch_{sfx}))")

        # 生成代码：添加公式标签文本框
        code.append(
            f"txBox = {slide_var}.shapes.add_textbox(current_col_x_{sfx}, y_pos, Inches(label_textbox_width_inch_{sfx}), Inches(element_height_inch_{sfx}))")
        code.append(
            "tf = txBox.text_frame; tf.clear(); tf.margin_bottom = Inches(0); tf.margin_top = Inches(0); tf.margin_left = Inches(0); tf.margin_right = Inches(0)")
        code.append("tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE")
        code.append("p = tf.paragraphs[0]")
        code.append(f"p.text = '{equ_text_label}'")
        code.append("p.font.size = Pt(16)")
        code.append("p.alignment = PP_ALIGN.LEFT")

        # 生成代码：添加公式图片
        code.append(f"img_x_pos = current_col_x_{sfx} + Inches(label_textbox_width_inch_{sfx})")
        code.append(
            f"pic = {slide_var}.shapes.add_picture(r'{equ_img_path}', img_x_pos, y_pos, "
            f"width=Inches(element_width_inch_{sfx}))")

        # 生成代码：更新当前列的最大宽度
        code.append(f"current_item_width = Inches(label_textbox_width_inch_{sfx}) + pic.width")
        code.append(f"if current_item_width > max_width_in_col_{sfx}: max_width_in_col_{sfx} = current_item_width")

        # 生成代码：更新行索引，并在需要时换列
        code.append(f"row_index_{sfx} += 1")
        code.append(f"if row_index_{sfx} >= max_rows_{sfx}:")
        code.append(f"    row_index_{sfx} = 0")
        code.append(f"    current_col_x_{sfx} += max_width_in_col_{sfx} + Inches(horizontal_gap_inch_{sfx})")
        code.append(f"    max_width_in_col_{sfx} = 0")

    code.append(f"# --- 结束公式网格布局 (后缀: {sfx}) ---\n")
    return code

def sanitize_for_var(name):
    # Convert any character that is not alphanumeric or underscore into underscore.
    return re.sub(r'[^0-9a-zA-Z_]+', '_', name)


def initialize_poster_code(width, height, slide_object_name, presentation_object_name, utils_functions):
    code = utils_functions
    code += fr'''
# Poster: {presentation_object_name}
{presentation_object_name} = create_poster(width_inch={width}, height_inch={height})

# Slide: {slide_object_name}
{slide_object_name} = add_blank_slide({presentation_object_name})
'''

    return code


def save_poster_code(output_file, utils_functions, presentation_object_name):
    code = utils_functions
    code = fr'''
# Save the presentation
save_presentation({presentation_object_name}, file_name="{output_file}")
'''
    return code


def generate_panel_code(panel_dict, utils_functions, slide_object_name, visible=False, theme=None):
    code = utils_functions
    raw_name = panel_dict["panel_name"]
    var_name = 'var_' + sanitize_for_var(raw_name)

    code += fr'''
# Panel: {raw_name}
{var_name} = add_textbox(
    {slide_object_name}, 
    '{var_name}', 
    {panel_dict['x']}, 
    {panel_dict['y']}, 
    {panel_dict['width']}, 
    {panel_dict['height']}, 
    text="", 
    word_wrap=True,
    font_size=40,
    bold=False,
    italic=False,
    alignment="left",
    fill_color=None,
    font_name="Arial"
)
'''

    if visible:
        if theme is None:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color=(0, 0, 0), thickness=5, line_style="solid")
'''
        else:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color={theme['color']}, thickness={theme['thickness']}, line_style="{theme['line_style']}")
'''

    return code


def generate_textbox_code(
        text_dict,
        utils_functions,
        slide_object_name,
        visible=False,
        content=None,
        theme=None,
        tmp_dir='tmp',
):
    code = utils_functions
    raw_name = text_dict["textbox_name"]
    var_name = sanitize_for_var(raw_name)

    code += fr'''
# Textbox: {raw_name}
{var_name} = add_textbox(
    {slide_object_name}, 
    '{var_name}', 
    {text_dict['x']}, 
    {text_dict['y']}, 
    {text_dict['width']}, 
    {text_dict['height']}, 
    text="", 
    word_wrap=True,
    font_size=40,
    bold=False,
    italic=False,
    alignment="left",
    fill_color=None,
    font_name="Arial"
)
'''
    if visible:
        if theme is None:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color=(255, 0, 0), thickness=5, line_style="solid")
'''
        else:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color={theme['color']}, thickness={theme['thickness']}, line_style="{theme['line_style']}")
'''

    if content is not None:
        tmp_name = f'{tmp_dir}/{var_name}_content.json'
        json.dump(content, open(tmp_name, 'w'), indent=4)
        code += fr'''
fill_textframe({var_name}, json.load(open('{tmp_name}', 'r')))
'''

    return code


def generate_figure_code(figure_dict, utils_functions, slide_object_name, img_path, visible=False, theme=None):
    code = utils_functions
    raw_name = figure_dict["figure_name"]
    var_name = sanitize_for_var(raw_name)

    code += fr'''
# Figure: {raw_name}
{var_name} = add_image(
    {slide_object_name}, 
    '{var_name}', 
    {figure_dict['x']}, 
    {figure_dict['y']}, 
    {figure_dict['width']}, 
    {figure_dict['height']}, 
    image_path="{img_path}"
)
'''

    if visible:
        if theme is None:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color=(0, 0, 255), thickness=5, line_style="long_dash_dot")
'''
        else:
            code += fr'''
# Make border visible
style_shape_border({var_name}, color={theme['color']}, thickness={theme['thickness']}, line_style="{theme['line_style']}")
'''

    return code





# ==============================================================================
# MODIFIED FUNCTIONS START HERE
# ==============================================================================

def _set_slide_background(prs, slide, image_path):
    """
    将一张图片设置为幻灯片的背景。
    这是一个处理 python-pptx 库限制的变通方案。
    """
    if not image_path or not os.path.exists(image_path):
        return

    # 将图片添加为覆盖整个幻灯片的大图
    pic = slide.shapes.add_picture(
        image_path,
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height
    )

    # 通过访问内部元素列表，将图片移动到最底层，使其成为背景
    # 这是目前社区公认的最可靠的方法
    slide.shapes._spTree.insert(2, pic._element)


def generate_multislide_ppt_code(raw_content,
                                 raw_result,
                                 equations,
                                 figure_arrangement_list,
                                 theme_config,
                                 save_path,
                                 slide_width_inch=16,
                                 slide_height_inch=9,
                                 title_bg_path='./assets/bg.jpg',
                                 content_bg_path='./assets/page.jpg'):
    """
    根据新的布局逻辑动态生成多页PPT的代码。
    """
    # --- 1. 初始化和数据准备 ---
    print('--------------- Running generate_multislide_ppt_code (Refactored Background Logic) -----------')

    from utils.theme_utils import get_slide_dimensions, get_theme_colors
    slide_width_inch, slide_height_inch = get_slide_dimensions(theme_config)
    title_text_color, content_style = get_theme_colors(theme_config)

    code_parts = [
        "import pptx",
        "import os",
        "from pptx.util import Inches, Pt",
        "from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR",
        "from pptx.dml.color import RGBColor",
        "import math",
        "",
        "def _set_slide_background(prs, slide, image_path):",
        "    if not image_path or not os.path.exists(image_path): return",
        "    pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)",
        "    slide.shapes._spTree.insert(2, pic._element)",
        "",
        "prs = pptx.Presentation()",
        f"prs.slide_width = Inches({slide_width_inch})",
        f"prs.slide_height = Inches({slide_height_inch})",
        "",
        "blank_layout = prs.slide_layouts[6]",
        "",
    ]

    figures_by_panel = {fig.get("slide_id") - 1: [] for fig in figure_arrangement_list if
                        fig.get("slide_id") is not None}
    for fig in figure_arrangement_list:
        if fig.get("slide_id") is not None:
            figures_by_panel[fig.get("slide_id") - 1].append(fig)

    ppt_page_content = {}
    im_id = 0

    # --- 2. 首页/封面页处理 ---
    if raw_content["sections"]:
        paper_title = "Paper Title"
        paper_authors = "Authors"
        paper_affiliations = ""

        if "meta" in raw_result:
            meta = raw_result["meta"]
            paper_title = meta.get("poster_title", meta.get("title", "Paper Title"))
            paper_authors = meta.get("authors", "Authors")
            paper_affiliations = meta.get("affiliations", "")

        slide_var = f"slide_0"
        code_parts.append(f"# ========= 开始生成封面页 =========\n")
        code_parts.append(f"{slide_var} = prs.slides.add_slide(blank_layout)")

        # --- 更优雅的调用方式 ---
        code_parts.append(f"_set_slide_background(prs, {slide_var}, r'{title_bg_path}')")

        # ... (封面页标题和作者代码保持不变) ...
        title_text_data = {"runs": [{"text": paper_title}]}
        title_textbox_var = f"title_textbox_0"
        code_parts.append(f"title_left = Inches(1.0)")
        code_parts.append(f"title_top = Inches(2.5)")
        code_parts.append(f"title_width = Inches({slide_width_inch - 2.0})")
        code_parts.append(f"title_height = Inches(2.0)")
        code_parts.append(
            f"{title_textbox_var} = {slide_var}.shapes.add_textbox(title_left, title_top, title_width, title_height)")
        code_parts.append(f"title_tf = {title_textbox_var}.text_frame; title_tf.clear(); title_tf.word_wrap = True")
        code_parts.extend(generate_paragraph_code("title_tf", title_text_data, override_font_size=48, is_title=True,
                                                  title_text_color=title_text_color, content_style=content_style))
        authors_text = f"{paper_authors}\n{paper_affiliations}" if paper_affiliations else paper_authors
        authors_text_data = {"runs": [{"text": authors_text}]}
        subtitle_textbox_var = f"subtitle_textbox_0"
        code_parts.append(f"subtitle_left = Inches(1.0)")
        code_parts.append(f"subtitle_top = Inches(5.5)")
        code_parts.append(f"subtitle_width = Inches({slide_width_inch - 2.0})")
        code_parts.append(f"subtitle_height = Inches(1.5)")
        code_parts.append(
            f"{subtitle_textbox_var} = {slide_var}.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)")
        code_parts.append(
            f"subtitle_tf = {subtitle_textbox_var}.text_frame; subtitle_tf.clear(); subtitle_tf.word_wrap = True")
        code_parts.extend(
            generate_paragraph_code("subtitle_tf", authors_text_data, override_font_size=26, is_title=False,
                                    title_text_color=title_text_color, content_style=content_style))

        ppt_page_content[im_id] = "Poster Title & Author"
        im_id += 1

    # --- 3. 内容页面处理 ---
    for i, section_data in enumerate(raw_content["sections"]):
        if i == 0: continue

        slide_var = f"slide_{i}"
        code_parts.append(f"# ========= 开始生成幻灯片 {i}: {section_data['title']} =========\n")

        title = section_data["title"]
        content = section_data["content"]
        ppt_page_content[im_id] = title
        title_text_data = {"runs": [{"text": title}]}
        text_data = {'runs': [{'text': content}]}

        code_parts.append(f"{slide_var} = prs.slides.add_slide(blank_layout)")

        # --- 更优雅的调用方式 ---
        code_parts.append(f"_set_slide_background(prs, {slide_var}, r'{content_bg_path}')")

        # ... (后续所有布局代码保持不变) ...
        title_textbox_var = f"title_textbox_{i}"
        code_parts.append(f"title_left = Inches(1.0)")
        code_parts.append(f"title_top = Inches(0.4)")
        code_parts.append(f"title_width = Inches({slide_width_inch - 2.0})")
        code_parts.append(f"title_height = Inches(1.0)")
        code_parts.append(
            f"{title_textbox_var} = {slide_var}.shapes.add_textbox(title_left, title_top, title_width, title_height)")
        code_parts.append(f"title_tf = {title_textbox_var}.text_frame; title_tf.clear(); title_tf.word_wrap = False")
        code_parts.extend(generate_paragraph_code("title_tf", title_text_data, override_font_size=36, is_title=True,
                                                  title_text_color=title_text_color, content_style=content_style))
        char_count = len(content)
        body_font_size = 18
        if char_count < 250:
            body_font_size = 24
        elif char_count < 550:
            body_font_size = 22
        elif char_count < 950:
            body_font_size = 20
        figures_for_this_slide = figures_by_panel.get(i - 1, [])
        if figures_for_this_slide:
            fig_data = figures_for_this_slide[0]
            fig_path = fig_data.get('figure_path', '').replace('\\', '/')
            fig_width = fig_data.get('width', 16)
            fig_height = fig_data.get('height', 9)
            if not fig_path or not os.path.exists(fig_path):
                figures_for_this_slide = []
            else:
                aspect_ratio = fig_width / fig_height
                aspect_ratio_threshold = 16 / 9
                left_margin_inch = 0.5
                right_margin_inch = 0.5
                top_content_margin_inch = 1.4
                bottom_margin_inch = 0.5
                content_gap_inch = 0.2
                if aspect_ratio > aspect_ratio_threshold:
                    code_parts.append(f"\n# --- 布局: 宽图 (幻灯片 {i}) ---")
                    available_height = slide_height_inch - top_content_margin_inch - bottom_margin_inch
                    text_height_inch = (available_height - content_gap_inch) / 2
                    body_textbox_var = f"body_textbox_{i}"
                    code_parts.append(f"body_left=Inches({left_margin_inch}); body_top=Inches({top_content_margin_inch}); body_width=Inches({slide_width_inch - left_margin_inch - right_margin_inch}); body_height=Inches({text_height_inch})")
                    code_parts.append(f"{body_textbox_var} = {slide_var}.shapes.add_textbox(body_left, body_top, body_width, body_height)")
                    code_parts.append(f"body_tf = {body_textbox_var}.text_frame; body_tf.word_wrap = True")
                    code_parts.extend(generate_paragraph_code("body_tf", text_data, override_font_size=body_font_size,
                                                              title_text_color=title_text_color,
                                                              content_style=content_style))
                    img_top_inch = top_content_margin_inch + text_height_inch + content_gap_inch
                    img_width_inch = slide_width_inch - left_margin_inch - right_margin_inch
                    # ==================== START: REVISED AND ROBUST CODE ====================
                    # 1. 定义图片可以放置的“盒子”的位置和尺寸 (即内容区的下半部分)
                    img_box_left_inch = left_margin_inch
                    img_box_top_inch = img_top_inch
                    img_box_width_inch = slide_width_inch - left_margin_inch - right_margin_inch
                    img_box_height_inch = text_height_inch

                    # 2. 根据“盒子”和图片的宽高比，计算能使其恰好容纳下的最终尺寸
                    final_img_width_inch = 0
                    final_img_height_inch = 0

                    # 比较盒子与图片的宽高比，确定缩放基准
                    if (img_box_width_inch / img_box_height_inch) > aspect_ratio:
                        # 情况A: 盒子的宽高比 > 图片的宽高比 (盒子相对更“宽”)
                        # 说明图片的高度是限制因素，应让图片高度撑满盒子
                        final_img_height_inch = img_box_height_inch
                        final_img_width_inch = img_box_height_inch * aspect_ratio
                    else:
                        # 情况B: 盒子的宽高比 <= 图片的宽高比 (盒子相对更“高”或一样)
                        # 说明图片的宽度是限制因素，应让图片宽度撑满盒子
                        final_img_width_inch = img_box_width_inch
                        final_img_height_inch = img_box_width_inch / aspect_ratio

                    # 3. 计算图片左上角的最终位置，使其在“盒子”内水平居中、垂直置顶
                    final_img_left_inch = img_box_left_inch + (img_box_width_inch - final_img_width_inch) / 2.0
                    final_img_top_inch = img_box_top_inch  # 垂直方向上与盒子顶部对齐

                    # 4. 使用计算出的精确位置和尺寸添加图片
                    # 同时提供 width 和 height 是最明确无误的方式，能避免库的内部缩放逻辑产生意外结果
                    code_parts.append(
                        f"pic = {slide_var}.shapes.add_picture(r'{fig_path}', Inches({final_img_left_inch}), Inches({final_img_top_inch}), width=Inches({final_img_width_inch}), height=Inches({final_img_height_inch}))")
                    # ===================== END: REVISED AND ROBUST CODE =====================

                else:
                    code_parts.append(f"\n# --- 布局: 高图 (幻灯片 {i}) ---")
                    img_top_inch = top_content_margin_inch * 1.5
                    img_height_inch = ( slide_height_inch - top_content_margin_inch - bottom_margin_inch) * 0.8
                    code_parts.append(
                        f"pic = {slide_var}.shapes.add_picture(r'{fig_path}', Inches(0), Inches({img_top_inch}), height=Inches({img_height_inch}))")
                    code_parts.append(f"pic_left_emu = prs.slide_width - pic.width - Inches({right_margin_inch})")
                    code_parts.append(f"pic.left = pic_left_emu")
                    body_textbox_var = f"body_textbox_{i}"
                    text_width_code = f"(pic_left_emu / 914400.0) - {left_margin_inch} - {content_gap_inch}"
                    code_parts.append(
                        f"body_left=Inches({left_margin_inch}); body_top=Inches({img_top_inch}); body_width=Inches({text_width_code}); body_height=Inches({img_height_inch})")
                    code_parts.append(
                        f"{body_textbox_var} = {slide_var}.shapes.add_textbox(body_left, body_top, body_width, body_height)")
                    code_parts.append(f"body_tf = {body_textbox_var}.text_frame; body_tf.word_wrap = True")
                    code_parts.extend(generate_paragraph_code("body_tf", text_data, override_font_size=body_font_size,
                                                              title_text_color=title_text_color,
                                                              content_style=content_style))
        if not figures_for_this_slide:
            equations_for_this_slide = [eq for eq in equations.values() if
                                        eq.get("section") == title and eq.get("equ_path") and os.path.exists(
                                            eq.get("equ_path"))]
            content_top_inch = 1.5
            bottom_margin_inch = 0.5
            if equations_for_this_slide:
                code_parts.append(f"\n# --- 布局: 文本 + 公式 (幻灯片 {i}) ---")
            else:
                code_parts.append(f"\n# --- 布局: 纯文本 (幻灯片 {i}) ---")
                body_height = slide_height_inch - content_top_inch - bottom_margin_inch
                body_textbox_var = f"body_textbox_{i}"
                code_parts.append(
                    f"body_left=Inches(1.0); body_top=Inches({content_top_inch}); body_width=Inches({slide_width_inch - 2.0}); body_height=Inches({body_height})")
                code_parts.append(
                    f"{body_textbox_var} = {slide_var}.shapes.add_textbox(body_left, body_top, body_width, body_height)")
                code_parts.append(f"body_tf = {body_textbox_var}.text_frame; body_tf.word_wrap = True")
                code_parts.extend(generate_paragraph_code("body_tf", text_data, override_font_size=body_font_size,
                                                          title_text_color=title_text_color,
                                                          content_style=content_style))
        im_id += 1
        code_parts.append("")

    # --- 4. 保存文件 ---
    code_parts.extend([
        "# ========= 保存PPT文件 =========",
        f"prs.save(r'{save_path}')",
        f"print(f'PPT已成功保存到: {{r'{save_path}'}}')",
    ])
    print(ppt_page_content)

    return "\n".join(code_parts), ppt_page_content


def generate_paragraph_code(text_frame_var, para_data, override_font_size=None, is_title=False, title_text_color=None, content_style=None):
    """
    一个辅助函数，为单个段落生成代码。
    新增 override_font_size 和 is_title 参数用于动态调整字体和样式。
    """
    code = []
    if not isinstance(para_data, dict):
        return code

    code.append(f"p = {text_frame_var}.add_paragraph()")
    p_level = para_data.get('level', 0)
    code.append(f"p.level = {p_level}")

    # 根据是否为标题设置对齐方式
    alignment = para_data.get('alignment', 'left').upper()
    if is_title:
        code.append(f"p.alignment = PP_ALIGN.CENTER")
    elif alignment == 'CENTER':
        code.append(f"p.alignment = PP_ALIGN.CENTER")
    else:
        code.append(f"p.alignment = PP_ALIGN.LEFT")

    for run_data in para_data.get('runs', []):
        text = run_data.get('text', '').replace("'", "\\'").replace('"', '\\"')
        is_bold = run_data.get('bold', False)
        # 标题强制加粗
        if is_title:
            is_bold = True

        font_size = override_font_size if override_font_size is not None else para_data.get('font_size')

        code.append(f"run = p.add_run()")
        code.append(f"run.text = r'''{text}'''")
        if is_bold:
            code.append(f"run.font.bold = True")
        if font_size:
            code.append(f"run.font.size = Pt({font_size})")
        
        # 根据is_title参数自动设置颜色
        if is_title and title_text_color:
            # 标题使用title_style的颜色
            code.append(f"try:")
            code.append(f"    run.font.color.rgb = RGBColor({title_text_color[0]}, {title_text_color[1]}, {title_text_color[2]})")
            code.append(f"except:")
            code.append(f"    pass")
        elif not is_title and content_style:
            # 内容使用content_style的颜色
            content_text_color = tuple(content_style.get('text_color', [0, 0, 0]))
            code.append(f"try:")
            code.append(f"    run.font.color.rgb = RGBColor({content_text_color[0]}, {content_text_color[1]}, {content_text_color[2]})")
            code.append(f"except:")
            code.append(f"    pass")
        else:
            # 默认黑色
            code.append(f"try:")
            code.append(f"    run.font.color.rgb = RGBColor(0, 0, 0)")
            code.append(f"except:")
            code.append(f"    pass")

    return code