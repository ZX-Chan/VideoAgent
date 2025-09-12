import argparse
import json
import os
import time
from utils.wei_utils import get_agent_config, scale_to_target_area
from PosterAgent.parse_raw import parse_raw, gen_image_and_table
from PosterAgent.my_pipeline_utils import run_filter_image_table, generate_narration_with_agent, generate_ppt_from_agent, synthesize_tts_with_openai, render_pptx_from_json
from PosterAgent.tree_split_layout import to_inches
from PosterAgent.gen_outline_layout import gen_outline_layout_v2
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

units_per_inch = 25

def analyze_reference_pptx(reference_path):
    """分析参考PPTX每页字数、图片数、表格数，返回统计summary。"""
    if not reference_path or not os.path.exists(reference_path):
        return None
    prs = Presentation(reference_path)
    stats = []
    for slide in prs.slides:
        text_len = sum(len(shape.text) for shape in slide.shapes if hasattr(shape, 'text'))
        img_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # 13: PICTURE
        table_count = sum(1 for shape in slide.shapes if shape.shape_type == 19)  # 19: TABLE
        stats.append({'text_len': text_len, 'img_count': img_count, 'table_count': table_count})
    return stats

def extract_plain_text(content):
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        texts = []
        for para in content:
            for run in para.get("runs", []):
                texts.append(run.get("text", ""))
        return " ".join(texts)
    return str(content)

def main():
    parser = argparse.ArgumentParser(description='Paper2Video Pipeline')
    parser.add_argument('--poster_path', type=str, required=True)
    parser.add_argument('--model_name_t', type=str, default='4o')
    parser.add_argument('--model_name_v', type=str, default='4o')
    parser.add_argument('--index', type=int, default=0)
    parser.add_argument('--poster_name', type=str, default=None)
    parser.add_argument('--tmp_dir', type=str, default='tmp')
    parser.add_argument('--ppt_width_inches', type=float, default=13.33, help='PPT page width in inches ')
    parser.add_argument('--ppt_height_inches', type=float, default=7.5, help='PPT page height in inches')
    parser.add_argument('--template_path', type=str, default=None, help='Optional: path to a PPTX template file')
    parser.add_argument('--reference_path', type=str, default=None, help='Optional: path to a reference PPTX file for layout/content guidance')
    parser.add_argument('--ablation_no_tree_layout', action='store_true', help='Ablation study: no tree layout')
    parser.add_argument('--ablation_no_commenter', action='store_true', help='Ablation study: no commenter')
    parser.add_argument('--ablation_no_example', action='store_true', help='Ablation study: no example')
    parser.add_argument('--no_blank_detection', action='store_true', help='When overflow is severe, try this option.')
    args = parser.parse_args()

    start_time = time.time()
    poster_name = args.poster_path.split('/')[-2].replace(' ', '_')
    if args.poster_name is None:
        args.poster_name = poster_name
    else:
        poster_name = args.poster_name
    out_dir = os.path.join(args.tmp_dir, poster_name)
    os.makedirs(out_dir, exist_ok=True)
    detail_log = {}

    tree_split_dir = 'tree_splits'
    contents_dir = 'contents'
    images_tables_dir = f'<{args.model_name_t}_{args.model_name_v}>_images_and_tables'
    os.makedirs(tree_split_dir, exist_ok=True)
    os.makedirs(contents_dir, exist_ok=True)
    os.makedirs(images_tables_dir, exist_ok=True)

    tts_output_dir = os.path.join(out_dir, 'tts')
    os.makedirs(tts_output_dir, exist_ok=True)

    agent_config_t = get_agent_config(args.model_name_t)
    agent_config_v = get_agent_config(args.model_name_v)
    meta_json_path = args.poster_path.replace('paper.pdf', 'meta.json')
    if args.ppt_width_inches is not None and args.ppt_height_inches is not None:
        ppt_width = args.ppt_width_inches * units_per_inch
        ppt_height = args.ppt_height_inches * units_per_inch
    elif os.path.exists(meta_json_path):
        meta_json = json.load(open(meta_json_path, 'r'))
        ppt_width = meta_json['width']
        ppt_height = meta_json['height']
    else:
        ppt_width = 48 * units_per_inch
        ppt_height = 36 * units_per_inch

    ppt_width_inches = to_inches(ppt_width, units_per_inch)
    ppt_height_inches = to_inches(ppt_height, units_per_inch)

    if ppt_width_inches > 56 or ppt_height_inches > 56:
        if ppt_width_inches >= ppt_height_inches:
            scale_factor = 56 / ppt_width_inches
        else:
            scale_factor = 56 / ppt_height_inches
        ppt_width_inches  *= scale_factor
        ppt_height_inches *= scale_factor
        ppt_width  = ppt_width_inches  * units_per_inch
        ppt_height = ppt_height_inches * units_per_inch

    print(f'PPT size: {ppt_width_inches} x {ppt_height_inches} inches')

    total_input_tokens_t, total_output_tokens_t = 0, 0
    total_input_tokens_v, total_output_tokens_v = 0, 0


    # Step 1: Parse the raw poster
    raw_content_path = os.path.join(contents_dir, f'<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_raw_content.json')
    if os.path.exists(raw_content_path):
        print(f'[Skip] Raw content exists: {raw_content_path}')
    else:
        input_token, output_token, raw_result = parse_raw(args, agent_config_t, version=2)
        total_input_tokens_t += input_token
        total_output_tokens_t += output_token
        _, _, images, tables = gen_image_and_table(args, raw_result)
        print(f'Parsing token consumption: {input_token} -> {output_token}')
        detail_log['parser_in_t'] = input_token
        detail_log['parser_out_t'] = output_token



    # Step 2: Filter unnecessary images and tables
    images_filtered_path = os.path.join(images_tables_dir, f'{args.poster_name}_images_filtered.json')
    tables_filtered_path = os.path.join(images_tables_dir, f'{args.poster_name}_tables_filtered.json')
    if os.path.exists(images_filtered_path) and os.path.exists(tables_filtered_path):
        print(f'[Skip] Filtered images/tables exist.')
    else:
        input_token, output_token = run_filter_image_table(args, agent_config_t)
        total_input_tokens_t += input_token
        total_output_tokens_t += output_token
        print(f'Filter figures token consumption: {input_token} -> {output_token}')
        detail_log['filter_in_t'] = input_token
        detail_log['filter_out_t'] = output_token



    # Step 3: Generate narration 
    narration_json_path = os.path.join(out_dir, f'{args.poster_name}_narration.json')
    if os.path.exists(narration_json_path):
        print(f'[Skip] Narration exists: {narration_json_path}')
        with open(narration_json_path, 'r', encoding='utf-8') as f:
            narration_json = json.load(f)
        narration_in_token = narration_out_token = 0
    else:
        with open(raw_content_path, 'r', encoding='utf-8') as f:
            raw_content = json.load(f)
        narration_json, narration_in_token, narration_out_token = generate_narration_with_agent(raw_content, agent_config_t)
        total_input_tokens_t += narration_in_token
        total_output_tokens_t += narration_out_token
        print(f'Narration token consumption: {narration_in_token} -> {narration_out_token}')
        detail_log['narration_in_t'] = narration_in_token
        detail_log['narration_out_t'] = narration_out_token
        with open(narration_json_path, 'w', encoding='utf-8') as f:
            json.dump(narration_json, f, ensure_ascii=False, indent=2)



    # Step 3.5: TTS
    tts_audio_log_path = os.path.join(tts_output_dir, 'tts_audio_files.json')
    if os.path.exists(tts_audio_log_path):
        print(f'[Skip] TTS audio exists: {tts_audio_log_path}')
        with open(tts_audio_log_path, 'r', encoding='utf-8') as f:
            tts_audio_files = json.load(f)
    else:
        openai_api_key = os.environ.get('OPENAI_API_KEY')
        if not openai_api_key:
            print('[TTS] OPENAI_API_KEY not set, skip TTS synthesis.')
            tts_audio_files = {}
        else:
            tts_audio_files = synthesize_tts_with_openai(narration_json, tts_output_dir, openai_api_key, voice="alloy", language="zh")
        with open(tts_audio_log_path, 'w', encoding='utf-8') as f:
            json.dump(tts_audio_files, f, ensure_ascii=False, indent=2)



    # Step 4: 用gen_outline_layout_v2生成大纲（每个panel对应一页PPT）
    reference_stats = analyze_reference_pptx(args.reference_path)
    outline_in_token, outline_out_token, panels, figures = gen_outline_layout_v2(args, agent_config_t)
    detail_log['outline_in_t'] = outline_in_token
    detail_log['outline_out_t'] = outline_out_token

    # Step 4.5: 布局推理与单位转换（完全复用new_pipeline）
    from PosterAgent.tree_split_layout import main_train, main_inference, get_arrangments_in_inches, split_textbox
    from PosterAgent.gen_poster_content import gen_bullet_point_content
    from utils.wei_utils import char_capacity, style_bullet_content
    # 读取过滤后的图片/表格
    with open(images_filtered_path, 'r', encoding='utf-8') as f:
        filtered_images = json.load(f)
    with open(tables_filtered_path, 'r', encoding='utf-8') as f:
        filtered_tables = json.load(f)
    # 布局推理
    panel_model_params, figure_model_params = main_train()
    panel_arrangement, figure_arrangement, text_arrangement = main_inference(
        panels,
        panel_model_params,
        figure_model_params,
        int(ppt_width),
        int(ppt_height),
        shrink_margin=3
    )
    text_arrangement_title = text_arrangement[0]
    text_arrangement = text_arrangement[1:]
    text_arrangement_title_top, text_arrangement_title_bottom = split_textbox(text_arrangement_title, 0.8)
    text_arrangement = [text_arrangement_title_top, text_arrangement_title_bottom] + text_arrangement
    for i in range(len(figure_arrangement)):
        panel_id = figure_arrangement[i]['panel_id']
        panel_section_name = panels[panel_id]['section_name']
        figure_info = figures[panel_section_name]
        if 'image' in figure_info:
            figure_id = figure_info['image']
            if not figure_id in filtered_images:
                figure_path = filtered_images[str(figure_id)]['image_path']
            else:
                figure_path = filtered_images[figure_id]['image_path']
        elif 'table' in figure_info:
            figure_id = figure_info['table']
            if not figure_id in filtered_tables:
                figure_path = filtered_tables[str(figure_id)]['table_path']
            else:
                figure_path = filtered_tables[figure_id]['table_path']
        figure_arrangement[i]['figure_path'] = figure_path
    for text_arrangement_item in text_arrangement:
        num_chars = char_capacity(
            bbox=(text_arrangement_item['x'], text_arrangement_item['y'], text_arrangement_item['height'], text_arrangement_item['width'])
        )
        text_arrangement_item['num_chars'] = num_chars
    width_inch, height_inch, panel_arrangement_inches, figure_arrangement_inches, text_arrangement_inches = get_arrangments_in_inches(
        ppt_width, ppt_height, panel_arrangement, figure_arrangement, text_arrangement, units_per_inch
    )
    # 保存tree_split结果，兼容gen_bullet_point_content依赖
    tree_split_results = {
        'poster_width': ppt_width,
        'poster_height': ppt_height,
        'poster_width_inches': ppt_width_inches,
        'poster_height_inches': ppt_height_inches,
        'panels': panels,
        'panel_arrangement': panel_arrangement,
        'figure_arrangement': figure_arrangement,
        'text_arrangement': text_arrangement,
        'panel_arrangement_inches': panel_arrangement_inches,
        'figure_arrangement_inches': figure_arrangement_inches,
        'text_arrangement_inches': text_arrangement_inches,
    }
    tree_split_path = os.path.join(tree_split_dir, f'<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_tree_split_{args.index}.json')
    with open(tree_split_path, 'w') as f:
        json.dump(tree_split_results, f, indent=4)
    # Step 5: 内容生成（对齐new_pipeline，使用generate_ppt_from_agent，确保每页有images字段）
    bullet_content_path = os.path.join(contents_dir, f'<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_bullet_point_content_{args.index}.json')
    if os.path.exists(bullet_content_path):
        bullet_content = json.load(open(bullet_content_path, 'r'))
        print(f'[Skip] Bullet content exists: {bullet_content_path}')
        input_token_t = output_token_t = input_token_v = output_token_v = 0
    else:
        # 读取raw_content
        raw_content_path = os.path.join(contents_dir, f'<{args.model_name_t}_{args.model_name_v}>_{args.poster_name}_raw_content.json')
        with open(raw_content_path, 'r', encoding='utf-8') as f:
            raw_content = json.load(f)
        # 生成结构化PPT内容（每页含images字段）
        from PosterAgent.my_pipeline_utils import generate_ppt_from_agent
        bullet_content, input_token_t, output_token_t = generate_ppt_from_agent(
            raw_content, filtered_images, filtered_tables, agent_config_t
        )
        with open(bullet_content_path, 'w', encoding='utf-8') as f:
            json.dump(bullet_content, f, ensure_ascii=False, indent=2)
        input_token_v = output_token_v = 0
    total_input_tokens_t += input_token_t
    total_output_tokens_t += output_token_t
    total_input_tokens_v += input_token_v
    total_output_tokens_v += output_token_v
    detail_log['content_in_t'] = input_token_t
    detail_log['content_out_t'] = output_token_t
    detail_log['content_in_v'] = input_token_v
    detail_log['content_out_v'] = output_token_v
    # Step 6: 样式
    theme_title_text_color = (255, 255, 255)
    theme_title_fill_color = (47, 85, 151)
    for k, v in bullet_content[0].items():
        style_bullet_content(v, theme_title_text_color, theme_title_fill_color)
    for i in range(1, len(bullet_content)):
        curr_content = bullet_content[i]
        style_bullet_content(curr_content['title'], theme_title_text_color, theme_title_fill_color)
    # Step 7: 渲染多页PPT
    prs = Presentation()
    prs.slide_width = Inches(ppt_width_inches)
    prs.slide_height = Inches(ppt_height_inches)
    pptx_path = os.path.join(out_dir, f'{args.poster_name}.pptx')
    for idx, panel in enumerate(panels):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 标题
        title = bullet_content[idx].get('title', panel.get('section_name', f'Section {idx+1}')) or ""
        title = extract_plain_text(title)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(ppt_width_inches-1), Inches(1.0))
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        # 正文（优先渲染textbox1和textbox2）
        text1 = bullet_content[idx].get('textbox1', '')
        text2 = bullet_content[idx].get('textbox2', '')
        text1 = extract_plain_text(text1)
        text2 = extract_plain_text(text2)
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(ppt_width_inches-1), Inches(3.5))
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        content_frame.clear()
        if text1:
            p = content_frame.add_paragraph()
            p.text = text1
            p.font.size = Pt(28)
            p.line_spacing = 1.3
        if text2:
            p2 = content_frame.add_paragraph()
            p2.text = text2
            p2.font.size = Pt(28)
            p2.line_spacing = 1.3
        # 图片（如有）
        img_list = []
        if 'images' in bullet_content[idx]:
            img_list = bullet_content[idx]['images']
        elif 'images' in panel:
            img_list = panel['images']
        img_list = [str(img) for img in img_list]
        for i, img_key in enumerate(img_list):
            img_info = filtered_images.get(img_key, None)
            if img_info and 'image_path' in img_info:
                img_path = img_info['image_path']
            else:
                continue
            if os.path.exists(img_path):
                try:
                    disp_w = ppt_width_inches * 0.6
                    disp_h = ppt_height_inches * 0.3
                    left = (ppt_width_inches - disp_w) / 2
                    top = ppt_height_inches - disp_h - 0.5
                    slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))
                except Exception as e:
                    print(f"[PPTX] Failed to add image {img_path}: {e}")
    prs.save(pptx_path)
    print(f"[PPTX] Saved to {pptx_path}")





    # 日志与耗时
    end_time = time.time()
    time_taken = end_time - start_time
    log_file = os.path.join(out_dir, f'{args.poster_name}_log.json')
    with open(log_file, 'w', encoding='utf-8') as f:
        log_data = {
            'input_tokens_t': total_input_tokens_t,
            'output_tokens_t': total_output_tokens_t,
            'input_tokens_v': total_input_tokens_v,
            'output_tokens_v': total_output_tokens_v,
            'time_taken': time_taken,
            'tts_audio_files': tts_audio_files,
            'pptx_path': pptx_path
        }
        json.dump(log_data, f, indent=4)
    detail_log_file = os.path.join(out_dir, f'{args.poster_name}_detail_log.json')
    with open(detail_log_file, 'w', encoding='utf-8') as f:
        json.dump(detail_log, f, indent=4)
    print(f'\nTotal time: {time_taken:.2f} seconds')
    print(f'Output directory: {out_dir}')
    print(f'Narration JSON: {narration_json_path}')
    print(f'TTS audio files: {tts_audio_log_path}')
    print(f'PPT structure JSON: {tree_split_path}')
    print(f'PPTX file: {pptx_path}')

if __name__ == '__main__':
    main()
