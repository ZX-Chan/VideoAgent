import os
import shutil
import subprocess
import tempfile
import traceback
from time import sleep, time
from types import SimpleNamespace

import json_repair
import Levenshtein
from lxml import etree
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from rich import print
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed
import moviepy
from moviepy import (ImageClip, concatenate_videoclips, ImageSequenceClip,
                            CompositeVideoClip, VideoFileClip)
from moviepy import vfx
import imageio
IMAGE_EXTENSIONS = {"bmp", "jpg", "jpeg", "pgm", "png", "ppm", "tif", "tiff", "webp"}

BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str):
    if file.split(".")[-1].lower() in IMAGE_EXTENSIONS:
        return True
    return False


def get_font_pptcstyle(font: dict):
    font = SimpleNamespace(**font)
    return f"Font Style: bold={font.bold}, italic={font.italic}, underline={font.underline}, size={font.size}pt, color={font.color}, font style={font.name}\n"


def get_font_style(font: dict):
    font = SimpleNamespace(**font)
    styles = []
    if font.size:
        styles.append(f"font-size: {font.size}pt")
    if font.color:
        styles.append(f"color: #{font.color}")
    if font.bold:
        styles.append("font-weight: bold")
    if font.italic:
        styles.append("font-style: italic")
    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph):
    runs = paragraph.runs
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath, seconds: int = 10, wait: bool = False):
    if not os.path.exists(filepath):
        while wait:
            print("waiting for:", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str):
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def get_slide_content(doc_json: dict, slide_title: str, slide: dict):
    slide_desc = slide.get("description", "")
    slide_content = f"Slide Purpose: {slide_title}\nSlide Description: {slide_desc}\n"
    for key in slide.get("subsections", []):
        slide_content += "Slide Content Source: "
        for section in doc_json["sections"]:
            subsections = section.get("subsections", [])
            if isinstance(subsections, dict) and len(subsections) == 1:
                subsections = [
                    {"title": k, "content": v} for k, v in subsections.items()
                ]
            for subsection in subsections:
                try:
                    if edit_distance(key, subsection["title"]) > 0.9:
                        slide_content += f"# {key} \n{subsection['content']}\n"
                except:
                    pass
    return slide_content


def tenacity_log(retry_state: RetryCallState):
    print(retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(raw_response: str):
    response = raw_response.strip()
    l, r = response.rfind("```json"), response.rfind("```")
    try:
        if l == -1 or r == -1:
            response = json_repair.loads(response)
        else:
            response = json_repair.loads(response[l + 7 : r].strip())
        return response
    except Exception as e:
        raise RuntimeError("Failed to parse JSON from response", e)


tenacity = retry(
    wait=wait_fixed(3), stop=stop_after_attempt(5), after=tenacity_log, reraise=True
)


@tenacity
def ppt_to_images(file: str, output_dir: str, warning: bool = False, dpi=72, output_type='png'):
    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir) and warning:
        print(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                if output_type == 'png':
                    img.save(pjoin(output_dir, f"poster.png"), 'PNG')
                else:
                    img.save(pjoin(output_dir, f"poster.jpg"), 'JPEG')
            return

        raise RuntimeError("No PDF file was created in the temporary directory", file)



@tenacity
def ppt_to_multi_images(file: str, output_dir: str, warning: bool = False, dpi=72, output_type='png'):
    """
    将PPT或PPTX文件的每一页转换为单独的图片文件。

    :param file: 输入的PPT/PPTX文件路径。
    :param output_dir: 输出图片的目录。
    :param warning: 如果输出目录已存在，是否打印警告信息。
    :param dpi: 输出图片的DPI（每英寸点数）。
    :param output_type: 输出图片类型，'png' 或 'jpg'。
    """
    assert os.path.exists(file), f"文件 {file} 不存在"
    if os.path.exists(output_dir) and warning:
        print(f"ppt_to_multi_images: 输出目录 {output_dir} 已存在")

    os.makedirs(output_dir, exist_ok=True)

    # soffice (LibreOffice) 需要在系统中安装
    with tempfile.TemporaryDirectory() as temp_dir:
        # 步骤1: 使用 soffice 将 PPT 转换为 PDF
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]
        try:
            subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            print("错误：无法执行 'soffice' 命令。请确保 LibreOffice 或 OpenOffice 已安装并已添加到系统 PATH。")
            raise e

        # 步骤2: 遍历临时目录中的PDF文件
        pdf_found = False
        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue

            pdf_found = True
            temp_pdf = os.path.join(temp_dir, f)

            # 步骤3: 使用 pdf2image 将 PDF 的每一页转换为图片
            images = convert_from_path(temp_pdf, dpi=dpi)

            for i, img in enumerate(images):
                # 核心修改：使用页码 i 作为文件名
                file_name = f"{i}.{output_type}"
                output_path = os.path.join(output_dir, file_name)

                if output_type == 'png':
                    img.save(output_path, 'PNG')
                elif output_type == 'jpg':
                    # 对于JPG，背景可能需要处理为白色以避免黑色背景
                    img = img.convert('RGB')
                    img.save(output_path, 'JPEG')

            # 找到并处理完第一个PDF后即可退出循环
            break

        if not pdf_found:
            raise RuntimeError("在临时目录中没有创建PDF文件", file)

@tenacity
def wmf_to_images(blob: bytes, filepath: str):
    if not filepath.endswith(".jpg"):
        raise ValueError("filepath must end with .jpg")
    dirname = os.path.dirname(filepath)
    basename = os.path.basename(filepath).removesuffix(".jpg")
    with tempfile.TemporaryDirectory() as temp_dir:
        with open(pjoin(temp_dir, f"{basename}.wmf"), "wb") as f:
            f.write(blob)
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "jpg",
            pjoin(temp_dir, f"{basename}.wmf"),
            "--outdir",
            dirname,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

    assert pexists(filepath), f"File {filepath} does not exist"


def extract_fill(shape: BaseShape):
    if "fill" not in dir(shape):
        return None
    else:
        return shape.fill._xPr.xml


def apply_fill(shape: BaseShape, fill_xml: str):
    if fill_xml is None:
        return
    new_element = etree.fromstring(fill_xml)
    shape.fill._xPr.getparent().replace(shape.fill._xPr, new_element)


def parse_groupshape(groupshape: GroupShape):
    assert isinstance(groupshape, GroupShape)
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height
        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )
    return group_shape_xy


def is_primitive(obj):
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)
    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(obj, result=None, exclude=None):
    if result is None:
        result = {}
    exclude = DEFAULT_EXCLUDE.union(exclude or set())
    for attr in dir(obj):
        if attr in exclude:
            continue
        try:
            if not attr.startswith("_") and not callable(getattr(obj, attr)):
                attr_value = getattr(obj, attr)
                if "real" in dir(attr_value):
                    attr_value = attr_value.real
                if attr == "size" and isinstance(attr_value, int):
                    attr_value = Length(attr_value).pt

                if is_primitive(attr_value):
                    result[attr] = attr_value
        except:
            pass
    return result


def merge_dict(d1: dict, d2: list[dict]):
    if len(d2) == 0:
        return d1
    for key in list(d1.keys()):
        values = [d[key] for d in d2]
        if d1[key] is not None and len(values) != 1:
            values.append(d1[key])
        if values[0] is None or not all(value == values[0] for value in values):
            continue
        d1[key] = values[0]
        for d in d2:
            d[key] = None
    return d1


def images_to_video(image_dir: str,
                    durations: list,
                    output_video_path: str,
                    fps: int = 24,
                    overlay_gif_path: str = None): # 新增：可选的GIF文件路径参数
    """
    将一系列图片根据指定的时长合成为一个视频文件，并可选择在右下角叠加一个循环播放的GIF。

    该函数会按文件名的数字顺序（0.png, 1.png, 2.png, ...）读取图片。

    :param image_dir: 包含图片文件（如 0.png, 1.png ...）的目录。
    :param durations: 一个列表，包含每张图片应播放的秒数。列表长度应与图片数量一致。
    :param output_video_path: 输出视频的文件路径（例如 'output.mp4'）。
    :param fps: 视频的帧率 (Frames Per Second)。
    :param overlay_gif_path: (可选) 要叠加在右下角的GIF文件路径。
    """
    # 1. 获取所有图片文件并进行自然排序
    try:
        image_files = [f for f in os.listdir(image_dir) if f.endswith(('.png', '.jpg'))]
        image_files.sort(key=lambda x: int(os.path.splitext(x)[0]))
    except FileNotFoundError:
        print(f"错误：找不到目录 '{image_dir}'")
        return
    except ValueError:
        print(f"错误：目录 '{image_dir}' 中的文件名不符合纯数字命名规范（如 '0.png', '1.png'）。")
        return

    # 2. 验证图片数量和时长列表长度是否匹配
    if len(image_files) != len(durations):
        print(f"错误：图片数量 ({len(image_files)}) 与时长列表长度 ({len(durations)}) 不匹配。")
        return

    if not image_files:
        print("错误：在指定目录中没有找到任何图片文件。")
        return

    print(f"找到 {len(image_files)} 张图片，准备合成为视频...")

    # 3. 为每张图片创建一个设置了时长的视频剪辑
    clips = []
    for i, image_file in enumerate(image_files):
        image_path = os.path.join(image_dir, image_file)
        duration = durations[i]
        clip = ImageClip(image_path).set_duration(duration)
        clips.append(clip)

    # 4. 将所有剪辑拼接成一个主视频轨道
    main_clip = concatenate_videoclips(clips, method="compose")

    # --- 新增功能：叠加GIF ---
    final_clip = main_clip # 默认最终剪辑为主视频
    if overlay_gif_path:
        print(f"正在添加GIF动图叠加: {overlay_gif_path}")
        try:
            # a. 使用 imageio 读取 GIF 的所有帧
            print("使用 imageio 读取 GIF 帧...")
            gif_reader = imageio.get_reader(overlay_gif_path)
            meta_data = gif_reader.get_meta_data()
            gif_fps = meta_data.get('fps', fps)
            gif_frames = [frame for frame in gif_reader]
            print(gif_fps)

            # b. 使用 ImageSequenceClip 从帧列表创建视频剪辑
            gif_clip = ImageSequenceClip(gif_frames, fps=gif_fps)
            resized_gif = gif_clip.resize(0.6)

            # c. 设置GIF循环播放，时长与主视频一致
            looped_gif = resized_gif.fx(vfx.loop, duration=main_clip.duration)

            # d. 放置在右下角
            positioned_gif = looped_gif.set_position(("right", "bottom"))

            # e. 合并
            final_clip = CompositeVideoClip([main_clip, positioned_gif])

        except Exception as e:
            print(f"错误：无法加载或处理GIF文件 '{overlay_gif_path}'。错误信息: {e}")
            print("将继续创建不带GIF叠加的视频。")
            final_clip = main_clip
    # --- 新增功能结束 ---

    # 5. 写入视频文件
    try:
        final_clip.write_videofile(output_video_path, fps=fps, codec="libx264", threads=4, logger='bar')
        print(f"\n视频成功创建！已保存至：{output_video_path}")
    except Exception as e:
        print(f"\n写入视频文件时出错: {e}")

def dict_to_object(dict: dict, obj: object, exclude=None):
    if exclude is None:
        exclude = set()
    for key, value in dict.items():
        if key not in exclude:
            setattr(obj, key, value)


class Config:

    def __init__(self, rundir=None, session_id=None, debug=True):
        self.DEBUG = debug
        if session_id is not None:
            self.set_session(session_id)
        if rundir is not None:
            self.set_rundir(rundir)

    def set_session(self, session_id):
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str):
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")
        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool):
        self.DEBUG = debug

    def remove_rundir(self):
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)


pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename

if __name__ == "__main__":
    config = Config()
    print(config)
